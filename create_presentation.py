#!/usr/bin/env python3
"""
Script to generate a PowerPoint presentation from the SplitSmart project.
Requires python-pptx library: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed.")
    print("Install it with: pip install python-pptx")
    exit(1)

def create_presentation():
    """Create PowerPoint presentation for SplitSmart project."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    primary_color = RGBColor(0, 102, 204)  # Blue
    secondary_color = RGBColor(0, 153, 76)  # Green
    accent_color = RGBColor(204, 0, 0)  # Red
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SplitSmart: End-to-End Cryptography Solution"
    subtitle.text = "Protecting Data Against Eavesdropping, Modification, Spoofing, and Replay Attacks\n\nProject 2.7"
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Project Overview"
    tf = content.text_frame
    tf.text = "Objective: Design and implement end-to-end cryptography solution"
    p = tf.add_paragraph()
    p.text = "Application: SplitSmart - Secure expense splitting"
    p = tf.add_paragraph()
    p.text = "Use Case: Groups splitting expenses (roommates, friends, colleagues)"
    p = tf.add_paragraph()
    p.text = "Key Requirements:"
    p = tf.add_paragraph()
    p.text = "• Protect against 4 core attacks"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Use state-of-the-art cryptographic algorithms"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Implement using cryptography libraries"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Demonstrate effectiveness through working demos"
    p.level = 1
    
    # Slide 3: Application Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Application Architecture"
    tf = content.text_frame
    tf.text = "Client-Side:"
    p = tf.add_paragraph()
    p.text = "• User registration and authentication"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Expense submission interface"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Cryptographic operations (encryption, signing)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Session management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Server-Side:"
    p = tf.add_paragraph()
    p.text = "• User management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Blockchain ledger storage"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Message processing and verification"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Balance calculations"
    p.level = 1
    
    # Slide 4: Threat Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Threat Model - Four Core Attacks"
    tf = content.text_frame
    tf.text = "1. Data Eavesdropping"
    p = tf.add_paragraph()
    p.text = "• Threat: Attacker intercepts network traffic"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Risk: Unauthorized access to expense data"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Data Modification"
    p = tf.add_paragraph()
    p.text = "• Threat: Attacker modifies messages in transit"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Risk: Unauthorized changes to expense amounts"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Data Originator Spoofing"
    p = tf.add_paragraph()
    p.text = "• Threat: Attacker impersonates legitimate users"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Risk: Unauthorized expense submissions"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Data Replay"
    p = tf.add_paragraph()
    p.text = "• Threat: Attacker captures and replays old messages"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Risk: Duplicate expense entries"
    p.level = 1
    
    # Slide 5: Cryptographic Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Cryptographic Architecture - Three Layers"
    tf = content.text_frame
    tf.text = "Layer 1: Handshake-Level Authentication"
    p = tf.add_paragraph()
    p.text = "• Technology: Signed Diffie-Hellman Key Exchange"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Mutual authentication between client and server"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Establishes secure session with forward secrecy"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Layer 2: Per-Entry Authentication"
    p = tf.add_paragraph()
    p.text = "• Technology: RSA-PSS Digital Signatures"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Each expense signed by user's private key"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Provides non-repudiation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Layer 3: Per-Message Protection"
    p = tf.add_paragraph()
    p.text = "• Technology: Multiple AEAD Encryption Algorithms"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AES-256-GCM, ChaCha20-Poly1305, AES-256-CBC-HMAC"
    p.level = 1
    
    # Slide 6: Defense Against Eavesdropping
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Defense Against Attack #1: Eavesdropping"
    tf = content.text_frame
    tf.text = "Attack Scenario:"
    p = tf.add_paragraph()
    p.text = "• Attacker intercepts network traffic"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Attempts to read expense data"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Defense Mechanism:"
    p = tf.add_paragraph()
    p.text = "• End-to-End Encryption (AES-256-GCM / ChaCha20-Poly1305)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Session keys derived from Diffie-Hellman key exchange"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 256-bit key strength provides confidentiality"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result:"
    p = tf.add_paragraph()
    p.text = "✓ Attacker sees only ciphertext (random data)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Cannot decrypt without session key"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Confidentiality preserved"
    p.level = 1
    
    # Slide 7: Defense Against Modification
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Defense Against Attack #2: Modification"
    tf = content.text_frame
    tf.text = "Attack Scenario:"
    p = tf.add_paragraph()
    p.text = "• Attacker intercepts and modifies encrypted messages"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Attempts to change expense amounts or descriptions"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Defense Mechanism:"
    p = tf.add_paragraph()
    p.text = "• Authenticated Encryption (AEAD)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AES-256-GCM includes authentication tag"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Any modification breaks authentication tag"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result:"
    p = tf.add_paragraph()
    p.text = "✓ Modified messages detected immediately"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Authentication tag verification fails"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Modified messages rejected"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Integrity preserved"
    p.level = 1
    
    # Slide 8: Defense Against Spoofing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Defense Against Attack #3: Spoofing"
    tf = content.text_frame
    tf.text = "Attack Scenario:"
    p = tf.add_paragraph()
    p.text = "• Attacker tries to submit expenses as another user"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Attempts to impersonate legitimate users"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Defense Mechanism:"
    p = tf.add_paragraph()
    p.text = "• Digital Signatures (RSA-PSS)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Each expense signed with user's private key"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Server verifies signature using public key"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 2048-bit RSA provides ~112-bit security"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result:"
    p = tf.add_paragraph()
    p.text = "✓ Invalid signatures detected"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Spoofed messages rejected"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Only legitimate users can create expenses"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Authentication and non-repudiation provided"
    p.level = 1
    
    # Slide 9: Defense Against Replay
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Defense Against Attack #4: Replay"
    tf = content.text_frame
    tf.text = "Attack Scenario:"
    p = tf.add_paragraph()
    p.text = "• Attacker captures valid messages"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Replays them later to create duplicate entries"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Defense Mechanism:"
    p = tf.add_paragraph()
    p.text = "• Monotonic Counters"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Each user has a counter that increments with each message"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Server checks counter is strictly increasing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Old messages have lower counters and are rejected"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result:"
    p = tf.add_paragraph()
    p.text = "✓ Replayed messages detected"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Counter check prevents replay"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Old messages cannot be reused"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Replay protection achieved"
    p.level = 1
    
    # Slide 10: Blockchain Ledger
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Blockchain Ledger - Tamper Evidence"
    tf = content.text_frame
    tf.text = "Blockchain Structure:"
    p = tf.add_paragraph()
    p.text = "• Each expense stored as a block"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Blocks linked via cryptographic hashes"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Genesis block starts the chain"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Block Components:"
    p = tf.add_paragraph()
    p.text = "• Block height (sequential number)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Previous block hash"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Merkle root"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Block hash"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Entry data with signature"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result:"
    p = tf.add_paragraph()
    p.text = "✓ Any database modification detected"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Hash chain breaks on tampering"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Provides tamper-evident history"
    p.level = 1
    
    # Slide 11: Cryptographic Algorithms
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Cryptographic Algorithms Used"
    tf = content.text_frame
    tf.text = "Key Exchange:"
    p = tf.add_paragraph()
    p.text = "• Diffie-Hellman (2048-bit): Secure key exchange"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• HKDF-SHA256: Key derivation from shared secret"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• RSA-2048: Digital signatures and key exchange signatures"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Encryption:"
    p = tf.add_paragraph()
    p.text = "• AES-256-GCM: Authenticated encryption, hardware-accelerated"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• ChaCha20-Poly1305: Authenticated encryption, software-optimized"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AES-256-CBC-HMAC: Encrypt-then-MAC for compatibility"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Digital Signatures:"
    p = tf.add_paragraph()
    p.text = "• RSA-PSS: Probabilistic signature scheme"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 2048-bit keys: ~112-bit security level"
    p.level = 1
    
    # Slide 12: Implementation Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Implementation Details"
    tf = content.text_frame
    tf.text = "Technology Stack:"
    p = tf.add_paragraph()
    p.text = "• Language: Python 3.x"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Cryptography Library: cryptography (Python)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Web Framework: Flask"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Database: SQLite with blockchain structure"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Frontend: HTML, CSS, JavaScript"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Key Components:"
    p = tf.add_paragraph()
    p.text = "• Client Crypto: Encryption, signing, key exchange"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Server Crypto: Decryption, verification, session management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Ledger: Blockchain hash chain implementation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Storage: Database with blockchain schema"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Protocols: Message format definitions"
    p.level = 1
    
    # Slide 13: Demonstrations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Demonstration - Attack Scenarios"
    tf = content.text_frame
    tf.text = "Demo 1: Eavesdropping Attack"
    p = tf.add_paragraph()
    p.text = "• python demos/demo_eavesdropping.py"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Shows intercepted ciphertext"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Demonstrates inability to decrypt"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Demo 2: Modification Attack"
    p = tf.add_paragraph()
    p.text = "• python demos/demo_modification.py"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Shows message modification attempt"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Demonstrates tag verification failure"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Demo 3: Spoofing Attack"
    p = tf.add_paragraph()
    p.text = "• python demos/demo_spoofing.py"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Shows impersonation attempt"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Demonstrates signature verification failure"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Demo 4: Replay Attack"
    p = tf.add_paragraph()
    p.text = "• python demos/demo_replay.py"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Shows message replay attempt"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Demonstrates counter check failure"
    p.level = 1
    
    # Slide 14: Blockchain Tampering Demo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Demonstration - Blockchain Tampering"
    tf = content.text_frame
    tf.text = "Demo 5: Ledger Tampering"
    p = tf.add_paragraph()
    p.text = "• python demos/demo_tampering.py"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Process:"
    p = tf.add_paragraph()
    p.text = "1. Legitimate expenses added to blockchain"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Attacker modifies database directly"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Server detects tampering on restart"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Hash chain verification fails"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5. Tampering alert generated"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result:"
    p = tf.add_paragraph()
    p.text = "✓ Any modification to ledger detected"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Blockchain structure prevents undetected tampering"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Provides audit trail"
    p.level = 1
    
    # Slide 15: Security Properties
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Security Properties Achieved"
    tf = content.text_frame
    tf.text = "Confidentiality ✓"
    p = tf.add_paragraph()
    p.text = "• All messages encrypted end-to-end"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Multiple encryption algorithms"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 256-bit key strength"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Integrity ✓"
    p = tf.add_paragraph()
    p.text = "• Authentication tags detect modifications"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Hash chain detects tampering"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Digital signatures verify authenticity"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Authentication ✓"
    p = tf.add_paragraph()
    p.text = "• Mutual authentication during key exchange"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Per-entry signatures verify origin"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Cannot impersonate without private key"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Non-Repudiation ✓"
    p = tf.add_paragraph()
    p.text = "• Digital signatures provide proof"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Users cannot deny their expenses"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Replay Protection ✓"
    p = tf.add_paragraph()
    p.text = "• Monotonic counters prevent replay"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Old messages rejected"
    p.level = 1
    
    # Slide 16: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Conclusion"
    tf = content.text_frame
    tf.text = "Project Achievements:"
    p = tf.add_paragraph()
    p.text = "✓ Complete Implementation: Working end-to-end cryptography solution"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Attack Protection: All 4 attacks defended against"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ State-of-the-Art: Modern cryptographic algorithms"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Demonstrations: Working attack demos"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Documentation: Comprehensive documentation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Key Takeaways:"
    p = tf.add_paragraph()
    p.text = "• End-to-end encryption is essential for data protection"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Multiple layers provide defense in depth"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Blockchain provides tamper evidence"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Proper implementation requires careful design"
    p.level = 1
    
    # Slide 17: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Questions & Discussion"
    subtitle.text = "Thank You!"
    
    # Save presentation
    filename = "SplitSmart_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation created: {filename}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()


